import threading
import queue
import uuid
import magic  # python-magic for MIME type detection
import fitz  # PyMuPDF for PDF processing
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
from qdrant_client import QdrantClient
from typing import Any, List
from jobtools import VectorJob, JobRegister
import io

class MainProcessor(threading.Thread):
    """
    A background processor thread to handle tasks asynchronously,
    specifically for processing VectorJob instances with QdrantClient and FastEmbed.
    """

    def __init__(self, task_lock: threading.Lock, task_queue: queue.Queue, job_register: JobRegister):
        """
        Initializes the MainProcessor with necessary components.

        Args:
            task_lock (threading.Lock): A lock for synchronizing task access.
            task_queue (queue.Queue): The task queue.
            job_register (JobRegister): Register for managing jobs.
        """
        super(MainProcessor, self).__init__()
        self.task_lock = task_lock
        self.task_queue = task_queue
        self.job_register = job_register

        # Initialize Qdrant client with FastEmbed
        self.qdrant_client = QdrantClient("localhost", port=6333)

        # Initialize the python-magic MIME type detector
        self.mime = magic.Magic(mime=True)

    def run(self):
        """
        The main loop to process tasks from the task queue.
        """
        while True:
            job_uuid = self.task_queue.get()  # Block until a job is available
            job = self.job_register.get_job(job_uuid)
            if isinstance(job, VectorJob):
                job.set_status("processing")
                if job.get_task_type() == "upload":
                    self.process_upload(job)
                elif job.get_task_type() == "query":
                    self.process_query(job)
                job.set_status("completed")
            self.task_queue.task_done()

    def process_upload(self, job: VectorJob) -> None:
        """
        Process an upload job to add files to a vector collection.

        Args:
            job (VectorJob): The vector job for uploading files and metadata.
        """
        collection_name = job.get_collection()
        
        # Retrieve documents as byte content and metadata
        documents = job.get_files()  # List of binary content (bytes)
        metadata = [
            {
                "source": job.get_metadata().get("source"),
                "author": job.get_metadata().get("author")
            } for _ in documents
        ]
       
        # Extract text content based on file type using python-magic
        decoded_documents = [
            self.extract_text_from_file(content) for content in documents
        ]

        # Use Qdrant's `add` method to handle embedding and storage
        self.qdrant_client.add(
            collection_name=collection_name,
            documents=decoded_documents,
            metadata=metadata

        )

        job.set_result({"message": "Files uploaded successfully", "points_count": len(documents)})

    def process_query(self, job: VectorJob) -> None:
        """
        Process a query job to retrieve similar vectors from a collection.

        Args:
            job (VectorJob): The vector job for querying the vector store.
        """
        collection_name = job.get_collection()
        query_text = job.get_query()

        # Perform the search query using Qdrant's `query` method
        search_result = self.qdrant_client.query(
            collection_name=collection_name,
            query_text=query_text,
            limit=5,  # Number of results to retrieve
            with_payload=True
        )

        # Format the search results
        result_data = [
            {
                "score": point.score,
                "metadata": point.payload
            } for point in search_result
        ]

        job.set_result(result_data)

    def extract_text_from_file(self, file_content: bytes) -> str:
        """
        Extracts text from various file types, including PDF, DOCX, XLSX, PPTX, HTML, and default text files.

        Args:
            file_content (bytes): Binary content of the file.

        Returns:
            str: Extracted text content.
        """
        # Determine MIME type of the file using python-magic
        mime_type = self.mime.from_buffer(file_content)

        if mime_type == "application/pdf":
            return self.extract_text_from_pdf(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self.extract_text_from_docx(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return self.extract_text_from_xlsx(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self.extract_text_from_pptx(file_content)
        elif mime_type == "text/html":
            return self.extract_text_from_html(file_content)
        else:
            return file_content.decode("utf-8", errors="ignore")

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        pdf_text = ""
        pdf = fitz.open(stream=file_content, filetype="pdf")
        for page in pdf:
            pdf_text += page.get_text()
        pdf.close()
        return pdf_text

    def extract_text_from_docx(self, file_content: bytes) -> str:
        text = ""
        docx = Document(io.BytesIO(file_content))
        for para in docx.paragraphs:
            text += para.text + "\n"
        return text

    def extract_text_from_xlsx(self, file_content: bytes) -> str:
        text = ""
        workbook = load_workbook(io.BytesIO(file_content), data_only=True)
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) if cell else "" for cell in row])
                text += row_text + "\n"
        return text

    def extract_text_from_pptx(self, file_content: bytes) -> str:
        text = ""
        ppt = Presentation(io.BytesIO(file_content))
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def extract_text_from_html(self, file_content: bytes) -> str:
        soup = BeautifulSoup(file_content, "html.parser")
        return soup.get_text()
